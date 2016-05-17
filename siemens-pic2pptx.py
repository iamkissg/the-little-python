#!/usr/bin/env python3
# -*- coding: utf-8 -*-

'''自动下载西门子带ppt的页面的ppt, 并生成ppt
目前无法自动设置ppt自动切换时间'''

__author__ = 'Engine'


import re
import os
import pptx
import argparse
import pptx.util
import scipy.misc
from selenium import webdriver
from urllib.request import urlretrieve



# regex, matching urls of pic and start time of it
PIC_RESOURCES = re.compile(r'urls\[(\d+)\]=\"(\S+)\"\nsts.*\"(\d+)\"')


def comma_handler(str_with_comma):
    return tuple(str_with_comma.split(","))


def get_resources(CourseID):
    '''根据资源的url, 获取ppt图片'''
    baseUrl = "http://www.ad.siemens.com.cn"
    url = baseUrl + "/service/elearning/cn/Course.aspx?CourseID=" + CourseID

    # 动态网页, 用urllib.request.urlopen获取到的html会有问题
    # 此处用selenium来获取html文档
    driver = webdriver.PhantomJS(executable_path="/home/kissg/Tools/phantomjs/bin/phantomjs")
    driver.get(url)

    page_source = driver.page_source
    try:   # 查看页面上是否有ppt资源, 不能用if语句, 因为find_element_by_id失败会抛出异常
        driver.find_element_by_id("picture")
    except Exception:
        print("Sorry, the page doesn't contain PPT resource.")
        raise Exception

    pic_resources = re.findall(PIC_RESOURCES, page_source)  # 匹配图片资源的urls与start times
    with open(CourseID + ".csv", "w") as f:
        for elm in pic_resources:  # download the pic
            urlretrieve(baseUrl + elm[1], elm[2] + ".jpg")
            # 1. 序号, 路径, 时间
            f.write(elm[0] + ',' + elm[1] + ',' + elm[2] + '\n')


def make_pptx(CourseID):
    '''自动生成pptx文件'''
    with open(CourseID + '.csv', 'r') as f:
        pic_info = list(map(comma_handler, f.read().strip().split('\n')))
    prs = pptx.Presentation()  # new

    pic_width = int(prs.slide_width * 1.0)

    pt = 0  # 前置时间, 用于做减法
    # generate slide
    for elm in pic_info:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        print("第", elm[0], "张幻灯片持续时间为", int(elm[2]) - pt, "s")
        pt = int(elm[2])

        img = scipy.misc.imread(elm[2] + ".jpg")
        pic_height = int(pic_width * img.shape[0] / img.shape[1])
        slide.shapes.add_picture(elm[2] + ".jpg", 0, 0, pic_width, pic_height)  # add pic to slide

    prs.save("%s.pptx" % "semi-automation")


if __name__ == '__main__':
    arg_parser = argparse.ArgumentParser(description="For Yang's task.")
    arg_parser.add_argument("id",
                            help="CourseID is needed.")
    args = arg_parser.parse_args()
    if not os.path.exists(args.id + ".csv"):
        get_resources(args.id)
    make_pptx(args.id)
